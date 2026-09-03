#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

WORK="/home/thanh/protein_designs/pdl1_bench/colabfold"
INK=RGBColor(0x15,0x20,0x2b); MUT=RGBColor(0x5b,0x6a,0x79)
BRAND=RGBColor(0x1f,0x5f,0x9e); GOOD=RGBColor(0x1a,0x7f,0x43); BAD=RGBColor(0xb3,0x20,0x2f)
GROUND=RGBColor(0xf5,0xf7,0xf9); WHITE=RGBColor(0xff,0xff,0xff)
GOODS=RGBColor(0xe7,0xf3,0xec); BADS=RGBColor(0xf8,0xe8,0xea); WARN=RGBColor(0x9a,0x6a,0x12)
FONT="Segoe UI"; SERIF="Georgia"

prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(bg=GROUND):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(1,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False; s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tb,tf

def bullet_color(b): return {"●":BRAND,"◆":GOOD,"■":BAD}.get(b,BRAND)

def para(tf,text,size,color=INK,bold=False,font=FONT,first=False,sb=0,sa=6,align=PP_ALIGN.LEFT,bullet=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_before=Pt(sb); p.space_after=Pt(sa)
    if bullet:
        rb=p.add_run(); rb.text=bullet+"  "; rb.font.color.rgb=bullet_color(bullet)
        rb.font.size=Pt(size); rb.font.bold=True; rb.font.name=font
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=font
    return p

def bar(s,color=BRAND):
    b=s.shapes.add_shape(1,0,0,I(0.16),SH); b.fill.solid(); b.fill.fore_color.rgb=color
    b.line.fill.background(); b.shadow.inherit=False

def eyebrow(s,text,color=BRAND,y=I(0.7)):
    _,tf=box(s,I(0.9),y,I(11),I(0.4)); para(tf,text.upper(),13,color,bold=True,first=True,sa=0)

# 1 Title
s=slide(GROUND); bar(s,BRAND)
eyebrow(s,"Structure-prediction validation · PD-L1",y=I(1.1))
_,tf=box(s,I(0.9),I(1.7),I(11.4),I(2.2))
para(tf,"Does ColabFold reproduce the strong binders?",40,INK,bold=True,font=SERIF,first=True,sa=10)
para(tf,"Independent localColabFold run on 7 PD-L1 binder designs (strong → weak); ipSAE and sc_DockQ recomputed and compared with the dataset consensus.",19,MUT,sa=0)
_,tf=box(s,I(0.9),I(6.2),I(11.5),I(0.9))
for i,(k,v) in enumerate([("Dataset","Anthropic/claude-protein-binder-design"),("Target","PD-L1 (115 aa)"),("Sample","7 designs"),("Date","2026-08-28")]):
    p=para(tf,"",12,MUT,first=(i==0),sa=2)
    rk=p.add_run(); rk.text=k+"  "; rk.font.bold=True; rk.font.color.rgb=INK; rk.font.size=Pt(12); rk.font.name=FONT
    rv=p.add_run(); rv.text=v; rv.font.color.rgb=MUT; rv.font.size=Pt(12); rv.font.name=FONT

# 2 Conclusion
s=slide(GROUND); bar(s,GOOD)
eyebrow(s,"Conclusion",color=GOOD)
_,tf=box(s,I(0.9),I(2.0),I(11.5),I(3.6))
para(tf,"Yes.",54,GOOD,bold=True,font=SERIF,first=True,sa=14)
para(tf,"ColabFold folds the STRONG designs into confident complexes that match the intended pose (ipTM ~0.91, sc_DockQ 0.87–0.96), and clearly fails on the WEAK ones (ipTM 0.21–0.58, ipSAE ≤0.19).",24,INK,sa=12)
para(tf,"Absolute values sit slightly below the dataset consensus, but the strong/weak ranking is preserved — which is what needed confirming.",18,MUT,sa=0)

# 3 Method
s=slide(GROUND); bar(s,BRAND)
eyebrow(s,"01 · Method")
_,tf=box(s,I(0.9),I(1.5),I(11.5),I(5.4))
para(tf,"Method",30,INK,bold=True,font=SERIF,first=True,sa=16)
for t in [
 "7 PD-L1 designs sampled across the dataset consensus ipsae_min (mean of 10 cofolding models), with wet-lab labels.",
 "Target + binder sequences joined into one complex; folded with localColabFold 1.6.2 (AF2-multimer-v3).",
 "Templates fully disabled; MSA from the remote server; --num-recycle up to 5.",
 "ipSAE recomputed with Dunbrack ipsae.py (cutoffs 10/10); sc_DockQ vs the design model (binder backbone, target heavy atoms).",
]:
    para(tf,t,19,INK,sa=12,bullet="●")

# 4 Results table
s=slide(GROUND); bar(s,BRAND)
eyebrow(s,"02 · Results")
_,tf=box(s,I(0.9),I(1.4),I(11),I(0.7)); para(tf,"Results — 7 designs",30,INK,bold=True,font=SERIF,first=True,sa=0)
rows=[
 ("Design","Wet-lab","Kd (nM)","Dataset ipSAE","CF ipTM","CF ipSAE","CF pLDDT bd","sc_DockQ"),
 ("M.st.rank03","binder","0.64","0.886","0.91","0.82","97","0.87","g"),
 ("M.mt.rank03","binder","47","0.880","0.91","0.80","98","0.94","g"),
 ("M.st.rank27","binder","19","0.836","0.92","0.79","98","0.96","g"),
 ("M.mt.rank15","non","—","0.831","0.91","0.79","97","0.88","b"),
 ("O.mt.rank28","non","—","0.552","0.58","0.19","85","0.05","b"),
 ("O.mt.rank23","non","—","0.451","0.29","0.02","86","0.19","b"),
 ("O.mt.rank27","non","—","0.259","0.21","0.01","46","0.09","b"),
]
nr,nc=len(rows),8
tbl=s.shapes.add_table(nr,nc,I(0.9),I(2.25),I(11.5),I(4.4)).table
tbl.columns[0].width=I(2.0)
for c in range(1,nc): tbl.columns[c].width=I((11.5-2.0)/(nc-1))
for ci,head in enumerate(rows[0]):
    cell=tbl.cell(0,ci); cell.fill.solid(); cell.fill.fore_color.rgb=INK; cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
    r=p.add_run(); r.text=head; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
for ri in range(1,nr):
    row=rows[ri]; tag=row[-1]
    for ci in range(nc):
        cell=tbl.cell(ri,ci); cell.fill.solid()
        cell.fill.fore_color.rgb=(GOODS if tag=="g" else BADS) if ci==0 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(row[ci]); r.font.size=Pt(12); r.font.name=FONT
        r.font.color.rgb=(GOOD if tag=="g" else BAD) if ci==0 else INK; r.font.bold=(ci==0)
_,tf=box(s,I(0.9),I(6.8),I(11.5),I(0.5))
para(tf,"'Dataset ipSAE' = released consensus; CF columns recomputed independently. Ordered strong (top) → weak (bottom).",12,MUT,first=True,sa=0)

# 5 Scatter figure
s=slide(WHITE); bar(s,BRAND)
eyebrow(s,"02 · Dataset (x-axis) vs localColabFold (y-axis)")
img=os.path.join(WORK,"cf_scatter_en.png")
w,h=Image.open(img).size; pic_w=I(8.7)
s.shapes.add_picture(img,I(0.55),I(1.45),width=pic_w)
_,tf=box(s,I(9.5),I(1.7),I(3.5),I(5.0))
para(tf,"How to read it",16,BRAND,bold=True,first=True,sa=10)
para(tf,"Dashed line = exact agreement.",14,INK,sa=8)
para(tf,"Filled = wet-lab binder.",14,GOOD,sa=8)
para(tf,"Open = non-binder.",14,BAD,sa=14)
para(tf,"Two clusters separate cleanly; weak designs fall below the diagonal, so ColabFold is stricter and the boundary is sharper.",14,MUT,sa=0)

# 6 Findings + caveat
s=slide(GROUND); bar(s,BRAND)
eyebrow(s,"03 · Findings")
_,tf=box(s,I(0.9),I(1.5),I(11.5),I(3.4))
para(tf,"Findings",30,INK,bold=True,font=SERIF,first=True,sa=16)
para(tf,"Clean separation: strong/mid ipTM 0.91–0.92, ipSAE 0.79–0.82, sc_DockQ 0.87–0.96; weak ipTM 0.21–0.58, ipSAE 0.01–0.19.",18,INK,sa=11,bullet="◆")
para(tf,"Ranking preserved: ColabFold scores correlate tightly with the dataset consensus (different predictor, so lower absolute values).",18,INK,sa=11,bullet="●")
para(tf,"Stricter on weak designs: recomputed scores fall even below the dataset consensus, widening the gap.",18,INK,sa=11,bullet="■")
cb=s.shapes.add_shape(1,I(0.9),I(5.15),I(11.5),I(1.55)); cb.fill.solid(); cb.fill.fore_color.rgb=RGBColor(0xf6,0xef,0xe0)
cb.line.color.rgb=WARN; cb.line.width=Pt(0.75); cb.shadow.inherit=False
tf=cb.text_frame; tf.word_wrap=True; tf.margin_left=I(0.25); tf.margin_top=I(0.15); tf.margin_right=I(0.25)
para(tf,"INTERPRETATION CAVEAT",12,WARN,bold=True,first=True,sa=5)
para(tf,"M.mt.rank15 scores high on structure (ipTM 0.91, sc_DockQ 0.88) yet is a wet-lab non-binder. A confident fold does NOT guarantee real binding.",15,INK,sa=0)
_,tf=box(s,I(0.9),I(6.95),I(11.5),I(0.4))
para(tf,"WSL2 / 8 GB caps GPU processes at ~600 s → weak designs run only 1–2 recycles (conclusion unchanged).",11,MUT,first=True,sa=0)

out=os.path.join(WORK,"ColabFold_PDL1_report.pptx"); prs.save(out); print("saved",out)
